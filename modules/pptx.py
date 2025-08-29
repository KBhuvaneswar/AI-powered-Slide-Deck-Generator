from pptx import Presentation
from pptx.util import Pt

def create_pptx(slides: list[dict], output_file: str = "generated_slides.pptx"):
    """
    Create a PowerPoint file from slide outline.
    
    - First slide will be a title-only slide (no bullets).
    - Following slides will have title + bullet points.
    
    Parameters:
        slides (list[dict]): List of slides with 'title' and 'bullets'.
        output_file (str): Output filename (default: generated_slides.pptx).
    """
    prs = Presentation()

    for i, slide_data in enumerate(slides):
        if i == 0:
            # Use Title Slide layout (index 0) for the first slide
            slide_layout = prs.slide_layouts[0]
            slide = prs.slides.add_slide(slide_layout)

            # Set Title
            title_placeholder = slide.shapes.title
            title_placeholder.text = slide_data.get("title", "Untitled Presentation")

            # Optional: set subtitle if provided
            if "subtitle" in slide_data:
                subtitle_placeholder = slide.placeholders[1]
                subtitle_placeholder.text = slide_data["subtitle"]

        else:
            # Use Title and Content layout (index 1) for subsequent slides
            slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(slide_layout)

            # Set Title
            title_placeholder = slide.shapes.title
            title_placeholder.text = slide_data.get("title", f"Slide {i+1}")

            # Set Bullet Points
            content_placeholder = slide.placeholders[1]
            tf = content_placeholder.text_frame
            tf.clear()

            for bullet in slide_data.get("bullets", []):
                p = tf.add_paragraph()
                p.text = bullet
                p.level = 0  
                p.font.size = Pt(18)

    prs.save(output_file)
    return output_file
